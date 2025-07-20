# documents/ocr_utils.py

from PIL import Image
import pytesseract
import os
from pdf2image import convert_from_path
from docx import Document as DocxDocument
import pandas as pd
from pptx import Presentation
import fitz 
from .ai_summarizer import generate_summary


def extract_text(image_path):
    """
    Extract text from a given image or scanned PDF file using Tesseract OCR.
    """
    if not os.path.exists(image_path):
        return "Image path not found."

    try:
        text = pytesseract.image_to_string(Image.open(image_path))
        return text.strip()
    except Exception as e:
        return f"OCR error: {e}"

# Scanned PDF OCR
def extract_scanned_pdf_text(file_path):
    pages = convert_from_path(file_path)
    text = ''
    for page in pages:
        text += pytesseract.image_to_string(page)
    return text


# Word (.docx)
def extract_docx_text(file_path):
    try:
        doc = DocxDocument(file_path)
        text = '\n'.join([para.text for para in doc.paragraphs if para.text.strip() != ""])
        return text.strip() or "No readable text in Word document."
    except Exception as e:
        return f"Error reading Word file: {e}"
    
# Excel (.xlsx)
def extract_excel_text(file_path):
    df = pd.read_excel(file_path)
    return df.to_string()

# PowerPoint (.pptx)

def extract_pptx_text(file_path):
    prs = Presentation(file_path)
    text = ''
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + '\n'
    return text


# PDF (Text-based)
def extract_pdf_text(file_path):
    text = ''
    with fitz.open(file_path) as pdf:
        for page in pdf:
            text += page.get_text()
    return text
